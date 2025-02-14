import io
from fastapi import FastAPI, File, UploadFile
from docx import Document
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader
import phonenumbers
import fitz
import re

app = FastAPI()


def extract_text_from_docx(file):
    document = Document(file)
    all_text = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            all_text.append(paragraph.text)

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    all_text.append(cell.text)

    return re.sub(r"\s+", " ", "\n".join(all_text)).strip()


def extract_hyperlinks_with_text(pdf):
    doc = fitz.open(stream=pdf)
    hyperlinks_with_text = []

    for page_number in range(len(doc)):
        page = doc[page_number]
        links = page.get_links()

        for link in links:
            if "uri" in link:
                hyperlink = link["uri"]
                rect = link["from"]
                anchor_text = page.get_textbox(rect).strip()

                if anchor_text:
                    hyperlinks_with_text.append(f"({anchor_text})[{hyperlink}]")
                else:
                    hyperlinks_with_text.append(f"(Unknown Text)[{hyperlink}]")

    return hyperlinks_with_text


def extract_text_from_ppt(file):
    text = ""
    powerppt = Presentation(file)
    for _, slide in enumerate(powerppt.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text

    return text


@app.post("/extract")
async def extract_text_from_file(file: UploadFile = File(...)):
    if file.filename.endswith(".docx") or file.filename.endswith(".doc"):
        text = extract_text_from_docx(file.file)
        return {"text": text}
    elif file.filename.endswith(".pdf"):
        hyperlinks = extract_hyperlinks_with_text(io.BytesIO(file.file.read()))
        reader = PdfReader(file.file)
        text = " ".join(page.extract_text() for page in reader.pages)
        reader.close()

        text += "\n here are hyperlinks with anchors if they exist:-\n"
        for i in hyperlinks:
            text += i + "\n"

        return {"text": text}
    elif file.filename.endswith(".pptx") or file.filename.endswith(".ppt"):
        text = extract_text_from_ppt(file.file)
        return {"text": text}
    else:
        return {"error": "File type not supported"}


class PhoneNumberRequest(BaseModel):
    phone_number: str


@app.post("/normalize-phone-number")
async def normalize_phone_number(request: PhoneNumberRequest):
    phone_number = request.phone_number
    parsed_number = phonenumbers.parse(phone_number)
    formatted_number = f"+{parsed_number.country_code}{parsed_number.national_number}"
    return {"phone_number": formatted_number}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=11037, access_log=False)
